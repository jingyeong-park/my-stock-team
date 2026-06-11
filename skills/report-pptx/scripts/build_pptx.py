# -*- coding: utf-8 -*-
"""report-pptx 렌더러: spec JSON -> 디자인된 PPTX.

사용법: python build_pptx.py <spec.json> <output.pptx>

디자인 규칙(코드 고정):
- 16:9, KB 옐로우(#FFBC00) 포인트 + 그레이/화이트, 차분한 금융 리포트 톤
- 모든 텍스트 '맑은 고딕' (latin + East Asian 모두 지정 -> 글자 깨짐 방지)
- 표는 최대 행 수 초과 시 자동 절단 + "외 N행 생략" 표기 (슬라이드 밖 잘림 방지)
- 금지 표현(매수/매도 단정·목표가) 발견 시 exit code 2로 실패 -> 호출자가 수정 후 재실행
"""
import json
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

FONT = "맑은 고딕"
KB_YELLOW = RGBColor(0xFF, 0xBC, 0x00)
DARK = RGBColor(0x3A, 0x3A, 0x3A)
GRAY = RGBColor(0x8C, 0x8C, 0x8C)
LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MAX_TABLE_ROWS = 8       # 헤더 제외 본문 최대 행
MAX_CHART_POINTS = 30    # 차트 포인트 상한(초과 시 다운샘플)

FORBIDDEN = [
    "매수 추천", "매도 추천", "강력 매수", "강력 매도",
    "사야 합니다", "팔아야 합니다", "사야 한다", "팔아야 한다",
    "목표주가", "목표가 ", "적정주가",
]

DISCLAIMER = "본 자료는 학습용 분석이며 투자 권유가 아닙니다."


# ---------- 텍스트 헬퍼 ----------

def _set_run(run, text, size=12, bold=False, color=DARK):
    run.text = text
    f = run.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    # 한글 깨짐 방지: East Asian 폰트도 맑은 고딕으로 고정
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", FONT)


def _add_para(tf, text, size=12, bold=False, color=DARK, bullet=False,
              align=PP_ALIGN.LEFT, space_after=6, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    run = p.add_run()
    _set_run(run, ("•  " + text) if bullet else text, size, bold, color)
    return p


def _textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


# ---------- 공통 레이아웃 ----------

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, x, y, w, h, fill, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if not line:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _title_bar(slide, title):
    _rect(slide, Inches(0.55), Inches(0.42), Inches(0.13), Inches(0.46), KB_YELLOW)
    tf = _textbox(slide, Inches(0.85), Inches(0.3), Inches(11.5), Inches(0.7))
    _add_para(tf, title, size=22, bold=True, first=True, space_after=0)


def _footer(slide, page_no, total):
    tf = _textbox(slide, Inches(0.55), Inches(7.05), Inches(8), Inches(0.35))
    _add_para(tf, DISCLAIMER, size=8, color=GRAY, first=True, space_after=0)
    tf2 = _textbox(slide, Inches(12.2), Inches(7.05), Inches(0.9), Inches(0.35))
    _add_para(tf2, f"{page_no} / {total}", size=8, color=GRAY,
              align=PP_ALIGN.RIGHT, first=True, space_after=0)


# ---------- 표 ----------

def _fill_cell(cell, text, size=11, bold=False, color=DARK, fill=None,
               align=PP_ALIGN.LEFT):
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.03)
    cell.margin_bottom = Inches(0.03)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    _set_run(p.add_run(), text, size, bold, color)


def _table(slide, x, y, w, h, columns, rows):
    """헤더=KB옐로우, 본문=화이트/라이트그레이 줄무늬. 행 초과 시 절단."""
    omitted = 0
    if len(rows) > MAX_TABLE_ROWS:
        omitted = len(rows) - MAX_TABLE_ROWS
        rows = rows[:MAX_TABLE_ROWS]
    size = 11 if len(rows) <= 5 else 9
    shape = slide.shapes.add_table(len(rows) + 1, len(columns), x, y, w, h)
    tbl = shape.table
    for j, col in enumerate(columns):
        _fill_cell(tbl.cell(0, j), str(col), size=size, bold=True,
                   fill=KB_YELLOW, align=PP_ALIGN.CENTER)
    for i, row in enumerate(rows):
        band = WHITE if i % 2 == 0 else LIGHT
        for j in range(len(columns)):
            val = str(row[j]) if j < len(row) else ""
            _fill_cell(tbl.cell(i + 1, j), val, size=size, fill=band,
                       align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT)
    if omitted:
        tf = _textbox(slide, x, y + h + Inches(0.02), w, Inches(0.3))
        _add_para(tf, f"※ 지면 관계상 {omitted}행 생략", size=9, color=GRAY,
                  first=True, space_after=0)
    return tbl


# ---------- 슬라이드 빌더 ----------

def slide_cover(prs, spec):
    s = _blank(prs)
    _rect(s, 0, 0, Inches(0.35), SLIDE_H, KB_YELLOW)          # 좌측 포인트 밴드
    _rect(s, Inches(0.35), Inches(6.9), SLIDE_W - Inches(0.35), Inches(0.6), LIGHT)
    tf = _textbox(s, Inches(1.1), Inches(2.4), Inches(11), Inches(2.6))
    name = spec["ticker_name"]
    code = spec.get("ticker_code", "")
    _add_para(tf, f"{name} 리서치 리포트", size=40, bold=True, first=True, space_after=10)
    if code:
        _add_para(tf, f"종목코드 {code}", size=16, color=GRAY, space_after=4)
    _add_para(tf, f"작성일 {spec['date']}", size=16, color=GRAY, space_after=4)
    tf2 = _textbox(s, Inches(1.1), Inches(6.95), Inches(10), Inches(0.4))
    _add_para(tf2, f"Stock Team 협업 리서치  |  {DISCLAIMER}", size=10,
              color=GRAY, first=True, space_after=0)


def slide_overview(prs, spec, pg, total):
    s = _blank(prs)
    _title_bar(s, "종목 개요")
    tf = _textbox(s, Inches(0.85), Inches(1.3), Inches(11.6), Inches(5.4))
    for i, b in enumerate(spec["overview"]["bullets"][:7]):
        _add_para(tf, b, size=14, bullet=True, first=(i == 0), space_after=12)
    _footer(s, pg, total)


def slide_financials(prs, spec, pg, total):
    s = _blank(prs)
    _title_bar(s, "재무 요약")
    fin = spec["financials"]
    n_rows = min(len(fin["rows"]), MAX_TABLE_ROWS)
    tbl_h = Inches(0.42 * (n_rows + 1))
    _table(s, Inches(0.85), Inches(1.35), Inches(11.6), tbl_h,
           fin["columns"], fin["rows"])
    y = Inches(1.35) + tbl_h + Inches(0.45)
    tf = _textbox(s, Inches(0.85), y, Inches(11.6), Inches(2.2))
    for i, b in enumerate(fin.get("bullets", [])[:4]):
        _add_para(tf, b, size=12, bullet=True, first=(i == 0), space_after=8)
    if fin.get("note"):
        _add_para(tf, fin["note"], size=9, color=GRAY, space_after=0)
    _footer(s, pg, total)


def slide_price(prs, spec, pg, total):
    s = _blank(prs)
    _title_bar(s, "가격 / 추세")
    price = spec["price"]
    chart = price.get("chart")
    if chart and chart.get("categories") and chart.get("series"):
        cats = list(chart["categories"])
        sers = chart["series"]
        # 다운샘플: 포인트가 너무 많으면 균등 간격 추출 (마지막 포인트는 보존)
        if len(cats) > MAX_CHART_POINTS:
            step = max(1, len(cats) // MAX_CHART_POINTS)
            idx = list(range(0, len(cats), step))
            if idx[-1] != len(cats) - 1:
                idx.append(len(cats) - 1)
            cats = [cats[i] for i in idx]
            sers = [{"name": se["name"],
                     "values": [se["values"][i] for i in idx]} for se in sers]
        data = CategoryChartData()
        data.categories = cats
        for se in sers:
            data.add_series(se["name"], se["values"])
        gf = s.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(0.85), Inches(1.3),
                                Inches(7.3), Inches(4.6), data)
        ch = gf.chart
        ch.has_legend = len(sers) > 1
        ch.font.name = FONT
        ch.font.size = Pt(9)
        ch.series[0].format.line.color.rgb = KB_YELLOW
        ch.series[0].format.line.width = Pt(2.5)
        ch.series[0].smooth = False
        bx = Inches(8.45)
        bw = Inches(4.2)
    else:
        bx = Inches(0.85)
        bw = Inches(11.6)
    tf = _textbox(s, bx, Inches(1.35), bw, Inches(4.8))
    for i, b in enumerate(price.get("bullets", [])[:7]):
        _add_para(tf, b, size=11.5, bullet=True, first=(i == 0), space_after=9)
    if price.get("source"):
        tf2 = _textbox(s, Inches(0.85), Inches(6.45), Inches(11.6), Inches(0.4))
        _add_para(tf2, price["source"], size=9, color=GRAY, first=True, space_after=0)
    _footer(s, pg, total)


def slide_news(prs, spec, pg, total):
    s = _blank(prs)
    _title_bar(s, "뉴스 · 시장 심리")
    news = spec["news"]
    tf = _textbox(s, Inches(0.85), Inches(1.3), Inches(11.6), Inches(4.4))
    for i, item in enumerate(news["bullets"][:5]):
        _add_para(tf, item["text"], size=12.5, bullet=True, first=(i == 0), space_after=2)
        _add_para(tf, "    " + item.get("source", ""), size=9, color=GRAY, space_after=10)
    if news.get("sentiment"):
        box = _rect(s, Inches(0.85), Inches(5.9), Inches(11.6), Inches(0.7), KB_YELLOW)
        btf = box.text_frame
        btf.word_wrap = True
        btf.margin_left = Inches(0.2)
        p = btf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        _set_run(p.add_run(), "시장 심리  |  " + news["sentiment"], 13, True, DARK)
    _footer(s, pg, total)


def slide_risks(prs, spec, pg, total):
    s = _blank(prs)
    _title_bar(s, "리스크")
    y = 1.3
    for i, r in enumerate(spec["risks"][:3]):
        _rect(s, Inches(0.85), Inches(y + 0.04), Inches(0.32), Inches(0.32), KB_YELLOW)
        tfn = _textbox(s, Inches(0.85), Inches(y), Inches(0.32), Inches(0.4))
        _add_para(tfn, str(i + 1), size=13, bold=True, align=PP_ALIGN.CENTER,
                  first=True, space_after=0)
        tf = _textbox(s, Inches(1.35), Inches(y), Inches(11.1), Inches(1.8))
        _add_para(tf, r["title"], size=14, bold=True, first=True, space_after=3)
        _add_para(tf, "근거: " + r.get("basis", ""), size=11, color=DARK, space_after=2)
        _add_para(tf, "영향: " + r.get("impact", ""), size=11, color=GRAY, space_after=0)
        y += 1.85
    _footer(s, pg, total)


def slide_summary(prs, spec, pg, total):
    s = _blank(prs)
    _title_bar(s, "종합")
    summ = spec["summary"]
    _rect(s, Inches(1.6), Inches(3.05), Inches(1.2), Inches(0.09), KB_YELLOW)
    tf = _textbox(s, Inches(1.6), Inches(3.3), Inches(10.1), Inches(2.4))
    _add_para(tf, summ["line"], size=20, bold=True, first=True, space_after=14)
    for b in summ.get("bullets", [])[:3]:
        _add_para(tf, b, size=12, color=GRAY, bullet=True, space_after=6)
    _footer(s, pg, total)


# ---------- 가드레일 ----------

def _walk_strings(obj):
    if isinstance(obj, str):
        yield obj
    elif isinstance(obj, dict):
        for v in obj.values():
            yield from _walk_strings(v)
    elif isinstance(obj, list):
        for v in obj:
            yield from _walk_strings(v)


def check_forbidden(spec):
    hits = []
    for text in _walk_strings(spec):
        for pat in FORBIDDEN:
            if pat in text:
                hits.append((pat, text[:80]))
    return hits


# ---------- 메인 ----------

def main():
    # Windows cp949 콘솔에서 한글 메시지 깨짐 방지
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    if len(sys.argv) != 3:
        print("사용법: python build_pptx.py <spec.json> <output.pptx>")
        sys.exit(1)
    spec = json.loads(Path(sys.argv[1]).read_text(encoding="utf-8"))

    hits = check_forbidden(spec)
    if hits:
        print("[금지 표현 검출] 아래 내용을 수정한 뒤 다시 실행하세요:")
        for pat, ctx in hits:
            print(f"  - '{pat}' in: {ctx}...")
        sys.exit(2)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = 7
    slide_cover(prs, spec)
    slide_overview(prs, spec, 2, total)
    slide_financials(prs, spec, 3, total)
    slide_price(prs, spec, 4, total)
    slide_news(prs, spec, 5, total)
    slide_risks(prs, spec, 6, total)
    slide_summary(prs, spec, 7, total)

    out = Path(sys.argv[2])
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"OK: {out} ({out.stat().st_size:,} bytes)")


if __name__ == "__main__":
    main()
